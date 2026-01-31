from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, List
import os
import json
import io
import httpx
from dotenv import load_dotenv
from openai import OpenAI
from supabase import create_client, Client
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import PyPDF2
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

load_dotenv()

app = FastAPI(title="Research Assistant Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
supabase: Client = create_client(
    os.getenv("SUPABASE_URL"),
    os.getenv("SUPABASE_SERVICE_KEY")
)

SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
FOLDER_ID = os.getenv("GOOGLE_DRIVE_FOLDER_ID")

def get_drive_service():
    creds_json = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON")
    creds_dict = json.loads(creds_json)
    creds = service_account.Credentials.from_service_account_info(creds_dict, scopes=SCOPES)
    return build('drive', 'v3', credentials=creds)

def extract_text_from_pdf(content: bytes) -> str:
    reader = PyPDF2.PdfReader(io.BytesIO(content))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_xlsx(content: bytes) -> str:
    wb = load_workbook(io.BytesIO(content))
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) for cell in row if cell])
            if row_text:
                text += row_text + "\n"
    return text

def get_embedding(text: str) -> List[float]:
    response = openai_client.embeddings.create(
        model="text-embedding-3-small",
        input=text[:8000]
    )
    return response.data[0].embedding

class ChatRequest(BaseModel):
    message: str
    history: Optional[List[dict]] = []

class SyncRequest(BaseModel):
    force: Optional[bool] = False

@app.get("/")
async def root():
    return {"status": "ok", "message": "Research Assistant Backend"}

@app.get("/health")
async def health():
    return {"status": "healthy"}

@app.post("/sync")
async def sync_documents(request: SyncRequest):
    try:
        service = get_drive_service()
        query = f"'{FOLDER_ID}' in parents and trashed = false"
        results = service.files().list(
            q=query,
            fields="files(id, name, mimeType, modifiedTime)"
        ).execute()
        
        files = results.get('files', [])
        synced = []
        
        for file in files:
            file_id = file['id']
            file_name = file['name']
            mime_type = file['mimeType']
            
            if not request.force:
                existing = supabase.table("documents").select("id").eq("metadata->>file_id", file_id).execute()
                if existing.data:
                    continue
            
            text = ""
            
            if mime_type == 'application/vnd.google-apps.document':
                req = service.files().export_media(fileId=file_id, mimeType='text/plain')
                content = req.execute()
                text = content.decode('utf-8')
            elif mime_type == 'application/vnd.google-apps.spreadsheet':
                req = service.files().export_media(fileId=file_id, mimeType='text/csv')
                content = req.execute()
                text = content.decode('utf-8')
            elif mime_type == 'application/vnd.google-apps.presentation':
                req = service.files().export_media(fileId=file_id, mimeType='text/plain')
                content = req.execute()
                text = content.decode('utf-8')
            elif mime_type == 'application/pdf':
                req = service.files().get_media(fileId=file_id)
                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, req)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                text = extract_text_from_pdf(fh.getvalue())
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                req = service.files().get_media(fileId=file_id)
                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, req)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                text = extract_text_from_docx(fh.getvalue())
            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                req = service.files().get_media(fileId=file_id)
                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, req)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                text = extract_text_from_pptx(fh.getvalue())
            elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                req = service.files().get_media(fileId=file_id)
                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, req)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                text = extract_text_from_xlsx(fh.getvalue())
            
            if text:
                chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
                for i, chunk in enumerate(chunks):
                    embedding = get_embedding(chunk)
                    supabase.table("documents").insert({
                        "content": chunk,
                        "metadata": {
                            "file_id": file_id,
                            "file_name": file_name,
                            "mime_type": mime_type,
                            "chunk_index": i
                        },
                        "embedding": embedding
                    }).execute()
                synced.append(file_name)
        
        return {"status": "success", "synced_files": synced, "total": len(synced)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/chat")
async def chat(request: ChatRequest):
    try:
        query_embedding = get_embedding(request.message)
        
        result = supabase.rpc(
            "match_documents",
            {
                "query_embedding": query_embedding,
                "match_threshold": 0.5,
                "match_count": 5
            }
        ).execute()
        
        context = ""
        sources = []
        if result.data:
            for doc in result.data:
                context += f"\n---\n{doc['content']}\n"
                if doc['metadata'].get('file_name') not in sources:
                    sources.append(doc['metadata'].get('file_name'))
        
        messages = [
            {
                "role": "system",
                "content": f"""당신은 연구 어시스턴트입니다. 
주어진 문서 컨텍스트를 바탕으로 질문에 답변해주세요.
컨텍스트에 없는 내용은 "해당 정보가 문서에 없습니다"라고 답변하세요.

문서 컨텍스트:
{context}
"""
            }
        ]
        
        for h in request.history[-10:]:
            messages.append({"role": h["role"], "content": h["content"]})
        
        messages.append({"role": "user", "content": request.message})
        
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0.7,
            max_tokens=1000
        )
        
        return {
            "answer": response.choices[0].message.content,
            "sources": sources
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/documents")
async def list_documents():
    try:
        result = supabase.table("documents").select("metadata").execute()
        files = {}
        for doc in result.data:
            file_name = doc['metadata'].get('file_name')
            if file_name and file_name not in files:
                files[file_name] = doc['metadata']
        return {"documents": list(files.values())}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=int(os.getenv("PORT", 8000)))
